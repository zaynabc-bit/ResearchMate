import pdfplumber
import os
from typing import Optional


def ocr_pdf_via_vision(file_path: str) -> str:
    """Run macOS Vision OCR on a scanned PDF using PyMuPDF and PyObjC."""
    try:
        import fitz  # PyMuPDF
        import objc
        from Foundation import NSData
        from Vision import VNImageRequestHandler, VNRecognizeTextRequest
    except ImportError as e:
        print(f"Failed to import OCR dependencies: {e}")
        return ""

    extracted_pages = []
    
    try:
        print(f"Opening PDF for OCR: {file_path}")
        with fitz.open(file_path) as doc:
            print(f"Running Vision OCR on {len(doc)} pages...")
            for page_num in range(len(doc)):
                page = doc[page_num]
                # Render page to a high-res image (150 DPI is a good balance of speed/quality)
                pix = page.get_pixmap(dpi=150)
                img_data = pix.tobytes("png")
                
                # Create NSData and request handler
                data = NSData.dataWithBytes_length_(img_data, len(img_data))
                handler = VNImageRequestHandler.alloc().initWithData_options_(data, None)
                
                # Setup OCR request
                request = VNRecognizeTextRequest.alloc().init()
                request.setRecognitionLevel_(0)  # 0 = Accurate (uses Neural Engine on Apple Silicon)
                
                # Run OCR
                success, error = handler.performRequests_error_([request], None)
                if success:
                    results = request.results()
                    page_text = []
                    for observation in results:
                        candidates = observation.topCandidates_(1)
                        if candidates:
                            page_text.append(candidates[0].string())
                    
                    text_content = "\n".join(page_text)
                    extracted_pages.append(text_content)
                    print(f"  Page {page_num + 1}/{len(doc)}: extracted {len(text_content)} characters")
                else:
                    print(f"  Page {page_num + 1}/{len(doc)}: Vision OCR failed: {error}")
                    extracted_pages.append("")
    except Exception as e:
        print(f"Error during Vision OCR: {e}")
        
    return "\n\n".join(extracted_pages)


def extract_text_from_pdf(file_path: str) -> dict:
    """Extract text and metadata from a PDF file with automatic OCR fallback."""
    result = {
        "text": "",
        "page_count": 0,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }

    try:
        result["file_size"] = os.path.getsize(file_path)

        with pdfplumber.open(file_path) as pdf:
            result["page_count"] = len(pdf.pages)

            # Extract metadata if available
            meta = pdf.metadata or {}
            result["title"] = meta.get("Title") or meta.get("title")
            result["authors"] = meta.get("Author") or meta.get("author")

            # Extract text from all pages
            pages_text = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text.strip())

            result["text"] = "\n\n".join(pages_text)

            # Automatic OCR fallback if extracted text is empty/extremely short
            total_chars = len(result["text"].strip())
            # If less than 100 characters total, or average is less than 50 characters per page
            if result["page_count"] > 0 and (total_chars < 100 or (total_chars / result["page_count"]) < 50):
                print(f"Extracted only {total_chars} characters via pdfplumber. Triggering macOS Vision OCR fallback...")
                ocr_text = ocr_pdf_via_vision(file_path)
                if ocr_text.strip():
                    result["text"] = ocr_text
                    print(f"Vision OCR completed. Extracted {len(result['text'])} characters total.")

            # Try to extract abstract from first page / text
            if result["text"]:
                first_page = result["text"][:3000].lower()
                abstract_start = first_page.find("abstract")
                if abstract_start != -1:
                    # Get text after "abstract" heading
                    abstract_text = result["text"][abstract_start + 8:abstract_start + 1200]
                    abstract_text = abstract_text.strip().lstrip(":")
                    result["abstract"] = abstract_text[:800] if abstract_text else None

            # If no title from metadata, try to get from first line of text
            if not result["title"] and result["text"]:
                first_lines = result["text"][:1000].split("\n")
                for line in first_lines[:5]:
                    line = line.strip()
                    if len(line) > 10 and len(line) < 200:
                        result["title"] = line
                        break

    except Exception as e:
        print(f"PDF extraction error: {e}")

    return result


def get_filename_title(filename: str) -> str:
    """Convert a filename to a readable title."""
    name = os.path.splitext(filename)[0]
    name = name.replace("_", " ").replace("-", " ")
    return name.title()


def extract_text_from_docx(file_path: str) -> dict:
    """Extract text from a Word document (.docx)."""
    result = {
        "text": "",
        "page_count": 1,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }
    try:
        import docx
        result["file_size"] = os.path.getsize(file_path)
        doc = docx.Document(file_path)
        
        # Read properties if available
        try:
            props = doc.core_properties
            result["title"] = props.title
            result["authors"] = props.author
        except Exception:
            pass
            
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        
        # Read from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
                        
        result["text"] = "\n\n".join(paragraphs)
        
        # Estimate page count
        word_count = len(result["text"].split())
        result["page_count"] = max(1, (word_count // 500) + 1)
        
        # Fallback title if properties are empty
        if not result["title"] and paragraphs:
            result["title"] = paragraphs[0][:120]
            
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        
    return result


def extract_text_from_txt(file_path: str) -> dict:
    """Extract text from a plain text file (.txt, .md)."""
    result = {
        "text": "",
        "page_count": 1,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }
    try:
        result["file_size"] = os.path.getsize(file_path)
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            result["text"] = f.read()
            
        lines = [line.strip() for line in result["text"].split("\n") if line.strip()]
        if lines:
            first_line = lines[0]
            if first_line.startswith("#"):
                first_line = first_line.lstrip("#").strip()
            result["title"] = first_line[:120]
            
        char_count = len(result["text"])
        result["page_count"] = max(1, (char_count // 3000) + 1)
        
    except Exception as e:
        print(f"TXT extraction error: {e}")
        
    return result


def extract_text_from_pptx(file_path: str, paper_id: str = None) -> dict:
    """Extract text and per-slide images from a PowerPoint presentation (.pptx)."""
    result = {
        "text": "",
        "page_count": 0,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
        "slides": [],  # [{slide_num, text, images: [relative_url, ...]}]
    }
    try:
        from pptx import Presentation
        import zipfile
        import re
        import shutil

        result["file_size"] = os.path.getsize(file_path)
        prs = Presentation(file_path)

        result["page_count"] = max(1, len(prs.slides))

        try:
            result["title"] = prs.core_properties.title
            result["authors"] = prs.core_properties.author
        except Exception:
            pass

        # Set up image output dir if paper_id given
        slides_dir = None
        if paper_id:
            slides_dir = os.path.join(os.path.dirname(file_path), "slides", paper_id)
            os.makedirs(slides_dir, exist_ok=True)

        # 1. Extract ALL images via ZIP relationships (catches backgrounds and groups)
        slide_images_map = {}  # slide_num -> [url, ...]
        if slides_dir and zipfile.is_zipfile(file_path):
            try:
                with zipfile.ZipFile(file_path, 'r') as z:
                    names = z.namelist()
                    for n in names:
                        m_rel = re.match(r'ppt/slides/_rels/slide(\d+)\.xml\.rels', n)
                        if m_rel:
                            s_num = int(m_rel.group(1))
                            content = z.read(n).decode('utf-8', errors='ignore')
                            media_targets = re.findall(r'Target="\.\./media/([^"]+)"', content)
                            if media_targets:
                                if s_num not in slide_images_map:
                                    slide_images_map[s_num] = []
                                for idx, media_name in enumerate(media_targets):
                                    media_path = f'ppt/media/{media_name}'
                                    if media_path in names:
                                        ext = media_name.split('.')[-1]
                                        img_filename = f"slide_{s_num:03d}_img_{idx + 1}.{ext}"
                                        img_out_path = os.path.join(slides_dir, img_filename)
                                        with z.open(media_path) as source, open(img_out_path, "wb") as target:
                                            shutil.copyfileobj(source, target)
                                        slide_images_map[s_num].append(f"/uploads/slides/{paper_id}/{img_filename}")
            except Exception as ze:
                print(f"  [PPTX] ZIP image extraction error: {ze}")

        # 2. Extract text with python-pptx
        all_text_parts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            slide_text = "\n".join(slide_texts)
            slide_images = slide_images_map.get(slide_num, [])

            if slide_text or slide_images:
                result["slides"].append({
                    "slide_num": slide_num,
                    "text": slide_text,
                    "images": slide_images,
                })
            if slide_text:
                all_text_parts.append(f"--- Slide {slide_num} ---\n{slide_text}")

        result["text"] = "\n\n".join(all_text_parts)

        if not result["title"] and result["slides"]:
            first_text = result["slides"][0].get("text", "")
            if first_text:
                result["title"] = first_text.splitlines()[0][:120]

    except Exception as e:
        print(f"PPTX extraction error: {e}")

    return result


def extract_text_from_csv(file_path: str) -> dict:
    """Extract text from a CSV file, representing rows as structured key-value sentences."""
    result = {
        "text": "",
        "page_count": 1,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }
    try:
        import csv
        result["file_size"] = os.path.getsize(file_path)
        result["title"] = os.path.basename(file_path).split('_', 1)[-1]
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            rows = list(reader)
            
        if not rows:
            return result
            
        headers = [h.strip() for h in rows[0] if h is not None]
        text_parts = []
        
        for idx, row in enumerate(rows[1:], start=1):
            if not row or all(not str(val).strip() for val in row):
                continue
            row_items = []
            for col_idx, val in enumerate(row):
                col_name = headers[col_idx] if col_idx < len(headers) else f"Column{col_idx+1}"
                row_items.append(f"{col_name}: {str(val).strip()}")
            row_text = f"Row {idx}: " + " | ".join(row_items)
            text_parts.append(row_text)
            
        result["text"] = f"Columns: {', '.join(headers)}\n\n" + "\n".join(text_parts)
        result["page_count"] = max(1, (len(text_parts) // 50) + 1)
        
    except Exception as e:
        print(f"CSV extraction error: {e}")
        
    return result


def extract_text_from_xlsx(file_path: str) -> dict:
    """Extract text from Excel sheet, representing sheets and rows as structured key-value sentences."""
    result = {
        "text": "",
        "page_count": 1,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }
    try:
        import openpyxl
        result["file_size"] = os.path.getsize(file_path)
        result["title"] = os.path.basename(file_path).split('_', 1)[-1]
        
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_sheets_text = []
        total_rows = 0
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for r in sheet.iter_rows(values_only=True):
                rows.append(r)
            
            if not rows:
                continue
                
            headers = [str(h).strip() if h is not None else f"Column{col_idx+1}" for col_idx, h in enumerate(rows[0])]
            sheet_parts = [f"--- Sheet: {sheet_name} ---", f"Columns: {', '.join(headers)}"]
            
            for idx, row in enumerate(rows[1:], start=1):
                if not row or all(val is None or not str(val).strip() for val in row):
                    continue
                row_items = []
                for col_idx, val in enumerate(row):
                    val_str = str(val).strip() if val is not None else ""
                    col_name = headers[col_idx] if col_idx < len(headers) else f"Column{col_idx+1}"
                    row_items.append(f"{col_name}: {val_str}")
                row_text = f"Row {idx}: " + " | ".join(row_items)
                sheet_parts.append(row_text)
                total_rows += 1
                
            all_sheets_text.append("\n".join(sheet_parts))
            
        result["text"] = "\n\n".join(all_sheets_text)
        result["page_count"] = max(1, (total_rows // 50) + 1)
        
    except Exception as e:
        print(f"XLSX extraction error: {e}")
        
    return result


def extract_text_from_doc(file_path: str) -> dict:
    """Extract text from a legacy binary Word document (.doc) using antiword."""
    result = {
        "text": "",
        "page_count": 1,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }
    try:
        import subprocess
        result["file_size"] = os.path.getsize(file_path)
        
        process = subprocess.Popen(
            ["antiword", file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        stdout, stderr = process.communicate()
        if process.returncode == 0:
            result["text"] = stdout.decode("utf-8", errors="ignore")
        else:
            print(f"Antiword error (code {process.returncode}): {stderr.decode('utf-8')}")
            
        if result["text"]:
            word_count = len(result["text"].split())
            result["page_count"] = max(1, (word_count // 500) + 1)
            
            lines = [line.strip() for line in result["text"].split("\n") if line.strip()]
            if lines:
                result["title"] = lines[0][:120]
        
    except Exception as e:
        print(f"DOC extraction error: {e}")
        
    return result


def extract_text_from_rtf(file_path: str) -> dict:
    """Extract text from a Rich Text Format file (.rtf) using striprtf."""
    result = {
        "text": "",
        "page_count": 1,
        "title": None,
        "authors": None,
        "abstract": None,
        "file_size": 0,
    }
    try:
        from striprtf.striprtf import rtf_to_text
        result["file_size"] = os.path.getsize(file_path)
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            rtf_content = f.read()
            
        result["text"] = rtf_to_text(rtf_content)
        
        if result["text"]:
            word_count = len(result["text"].split())
            result["page_count"] = max(1, (word_count // 500) + 1)
            
            lines = [line.strip() for line in result["text"].split("\n") if line.strip()]
            if lines:
                result["title"] = lines[0][:120]
                
    except Exception as e:
        print(f"RTF extraction error: {e}")
        
    return result


def extract_text_from_file(file_path: str, paper_id: str = None) -> dict:
    """Extract text and metadata from any supported document format."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".doc":
        return extract_text_from_doc(file_path)
    elif ext == ".rtf":
        return extract_text_from_rtf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path, paper_id=paper_id)
    elif ext == ".csv":
        return extract_text_from_csv(file_path)
    elif ext in [".xlsx", ".xls"]:
        return extract_text_from_xlsx(file_path)
    elif ext in [".txt", ".md"]:
        return extract_text_from_txt(file_path)
    else:
        return extract_text_from_txt(file_path)

